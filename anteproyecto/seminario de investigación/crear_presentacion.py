#!/usr/bin/env python3
"""
Script para generar presentación PowerPoint del Protocolo de Investigación
Seminario de Investigación I - UTP
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def crear_presentacion():
    """Crea la presentación completa del protocolo de investigación"""
    
    # Crear presentación
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # === SLIDE 1: PORTADA ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fondo de color sutil
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 248, 250)
    
    # Título principal
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = 'Modelo DLBP para Optimización del\nBalanceo de Carcasa en la Industria Avícola'
    title_frame.word_wrap = True
    for paragraph in title_frame.paragraphs:
        paragraph.font.size = Pt(38)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(0, 51, 102)
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.space_before = Pt(6)
    
    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = 'Enfoque con Técnicas Metaheurísticas'
    subtitle_frame.paragraphs[0].font.size = Pt(26)
    subtitle_frame.paragraphs[0].font.italic = True
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(70, 130, 180)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Información del autor
    author_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(1.8))
    author_frame = author_box.text_frame
    author_frame.text = 'Daniel Castañeda'
    p = author_frame.add_paragraph()
    p.text = 'Maestría en Investigación de Operaciones y Estadística'
    p = author_frame.add_paragraph()
    p.text = 'Universidad Tecnológica de Pereira'
    p = author_frame.add_paragraph()
    p.text = 'Seminario de Investigación I - 2025'
    for paragraph in author_frame.paragraphs:
        paragraph.font.size = Pt(17)
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.color.rgb = RGBColor(64, 64, 64)
        paragraph.space_after = Pt(4)
    
    # === SLIDE 2: AGENDA ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 248, 250)
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = 'Agenda'
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    # Contenido
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(8), Inches(5))
    tf = content_box.text_frame
    items = [
        '1. Planteamiento del Problema',
        '2. Justificación y Relevancia',
        '3. Objetivos de la Investigación',
        '4. Marco Teórico: DLBP y Metaheurísticas',
        '5. Metodología Propuesta',
        '6. Cronograma y Recursos',
        '7. Resultados Esperados'
    ]
    
    for i, item in enumerate(items):
        if i == 0:
            tf.text = item
        else:
            p = tf.add_paragraph()
            p.text = item
        tf.paragraphs[i].font.size = Pt(22)
        tf.paragraphs[i].font.color.rgb = RGBColor(30, 30, 30)
        tf.paragraphs[i].space_after = Pt(14)
        if i > 0:
            tf.paragraphs[i].space_before = Pt(0)
    
    # === SLIDE 3: CONTEXTO ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 248, 250)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = 'Contexto: Industria Avícola Colombiana'
    title_frame.paragraphs[0].font.size = Pt(38)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(8), Inches(5.2))
    tf = content_box.text_frame
    tf.text = '✓ Importancia económica estratégica'
    p = tf.add_paragraph()
    p.text = '✓ Producción de múltiples coproductos por carcasa'
    p = tf.add_paragraph()
    p.text = ''
    p = tf.add_paragraph()
    p.text = '⚠️ Problema Central: Desbalance de Carcasa'
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(178, 34, 34)
    p = tf.add_paragraph()
    p.text = ''
    p = tf.add_paragraph()
    p.text = '   • Oferta fija vs. Demanda variable'
    p = tf.add_paragraph()
    p.text = '   • Ineficiencias operativas y sobrecostos'
    p = tf.add_paragraph()
    p.text = '   • Excedentes de productos de bajo valor'
    p = tf.add_paragraph()
    p.text = '   • Faltantes de productos premium'
    
    for i, para in enumerate(tf.paragraphs):
        if i <= 1:
            para.font.size = Pt(22)
        elif i <= 8 and i >= 5:
            para.font.size = Pt(20)
        para.font.color.rgb = RGBColor(30, 30, 30) if i != 3 else RGBColor(178, 34, 34)
        para.space_after = Pt(8)
    
    # === SLIDE 4: EL PROBLEMA ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 248, 250)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = 'Planteamiento del Problema'
    title_frame.paragraphs[0].font.size = Pt(38)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.5), Inches(5.2))
    tf = content_box.text_frame
    tf.text = 'Desbalance Estructural'
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    items = [
        '• Cada ave → Proporciones fijas de cortes',
        '• Mercado → Demanda heterogénea y estocástica',
        '',
        'Consecuencias Económicas',
        '• Sobrecostos operativos documentados',
        '• Violación de restricciones de bioseguridad',
        '• Costos de inventario y almacenamiento',
        '• Ventas a precios de liquidación'
    ]
    
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        if item.startswith('Consecuencias'):
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 51, 102)
        else:
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(30, 30, 30)
        p.space_after = Pt(6)
    
    # === SLIDE 5: PREGUNTA DE INVESTIGACIÓN ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 248, 255)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = 'Pregunta de Investigación'
    title_frame.paragraphs[0].font.size = Pt(38)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Caja con borde para la pregunta
    question_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(3.5))
    question_frame = question_box.text_frame
    question_frame.text = '¿Cómo puede un modelo DLBP resuelto con metaheurísticas mejorar la rentabilidad y eficiencia operativa al minimizar el desbalance entre oferta y demanda de coproductos avícolas?'
    question_frame.word_wrap = True
    for paragraph in question_frame.paragraphs:
        paragraph.font.size = Pt(28)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(25, 25, 112)
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.line_spacing = 1.3
    
    # === SLIDE 6: JUSTIFICACIÓN ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 248, 250)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = 'Justificación: Beneficios Documentados'
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5))
    tf = content_box.text_frame
    
    sections = [
        ('Beneficios Económicos', True),
        ('• Optimización de la rentabilidad global', False),
        ('• Reducción de costos operativos', False),
        ('', False),
        ('Caso Santa Marta (Solano-Blanco, 2022)', True),
        ('• Reducción 8.6% en costos totales', False),
        ('• Cumplimiento de restricciones biológicas', False),
        ('', False),
        ('Beneficios de Sostenibilidad', True),
        ('• Alineación con ODS 12 y 9', False),
        ('• Reducción de desperdicios (Marel)', False)
    ]
    
    tf.text = sections[0][0]
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0, 100, 0)
    
    for text, is_header in sections[1:]:
        p = tf.add_paragraph()
        p.text = text
        if is_header:
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 100, 0)
        else:
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(30, 30, 30)
        p.space_after = Pt(6)
    
    # === SLIDE 7: OBJETIVOS ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 248, 250)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = 'Objetivos'
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5))
    tf = content_box.text_frame
    
    sections = [
        ('Objetivo General', True),
        ('Desarrollar y validar un modelo DLBP con metaheurísticas para maximizar rentabilidad en la gestión de coproductos avícolas', False),
        ('', False),
        ('Objetivos Específicos', True),
        ('1. Formular modelo matemático DLBP', False),
        ('2. Implementar algoritmos (GA, TS, Híbrido)', False),
        ('3. Generar dataset sintético calibrado', False),
        ('4. Evaluar desempeño y cuantificar impacto', False)
    ]
    
    tf.text = sections[0][0]
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    for text, is_header in sections[1:]:
        p = tf.add_paragraph()
        p.text = text
        if is_header:
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 51, 102)
        else:
            p.font.size = Pt(19)
            p.font.color.rgb = RGBColor(30, 30, 30)
        p.space_after = Pt(8)
    
    # === SLIDE 8: MARCO TEÓRICO - DLBP ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 248, 250)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = 'Marco Teórico: DLBP'
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5))
    tf = content_box.text_frame
    
    sections = [
        ('Disassembly Line Balancing Problem', True),
        ('• Proceso inverso al ensamblaje (uno → muchos)', False),
        ('• Optimización de líneas de desensamble', False),
        ('• Complejidad: NP-hard', False),
        ('', False),
        ('Aplicación Avícola', True),
        ('• Carcasa completa → Múltiples coproductos', False),
        ('• Restricciones de precedencia en despiece', False),
        ('• Balance de proporciones fijas', False),
        ('• Demanda estocástica', False)
    ]
    
    tf.text = sections[0][0]
    tf.paragraphs[0].font.size = Pt(23)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(70, 130, 180)
    
    for text, is_header in sections[1:]:
        p = tf.add_paragraph()
        p.text = text
        if is_header:
            p.font.size = Pt(23)
            p.font.bold = True
            p.font.color.rgb = RGBColor(70, 130, 180)
        else:
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(30, 30, 30)
        p.space_after = Pt(6)
    
    # === SLIDE 9: METAHEURÍSTICAS ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 248, 250)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = 'Técnicas Metaheurísticas Propuestas'
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5))
    tf = content_box.text_frame
    
    sections = [
        ('Algoritmos Genéticos (GA)', True),
        ('• Exploración global del espacio de soluciones', False),
        ('• Operadores: selección, cruce, mutación', False),
        ('', False),
        ('Búsqueda Tabú (TS)', True),
        ('• Búsqueda local con memoria', False),
        ('• Escape de óptimos locales', False),
        ('', False),
        ('Algoritmo Híbrido (GA-TS)', True),
        ('• Combina exploración (GA) + explotación (TS)', False),
        ('• Balance diversificación e intensificación', False)
    ]
    
    tf.text = sections[0][0]
    tf.paragraphs[0].font.size = Pt(23)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(70, 130, 180)
    
    for text, is_header in sections[1:]:
        p = tf.add_paragraph()
        p.text = text
        if is_header:
            p.font.size = Pt(23)
            p.font.bold = True
            p.font.color.rgb = RGBColor(70, 130, 180)
        else:
            p.font.size = Pt(19)
            p.font.color.rgb = RGBColor(30, 30, 30)
        p.space_after = Pt(6)
    
    # === SLIDE 10: METODOLOGÍA ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 248, 250)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = 'Metodología: 5 Fases'
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5))
    tf = content_box.text_frame
    
    phases = [
        ('Fase 1: Formulación del Modelo (Sem 1-8)', 'Revisión literatura + Formulación matemática'),
        ('Fase 2: Implementación (Sem 9-16)', 'Codificación GA, TS y Híbrido en Python'),
        ('Fase 3: Generación de Datos (Sem 17-20)', 'Dataset sintético calibrado'),
        ('Fase 4: Experimentación (Sem 21-24)', 'Pruebas + Análisis estadístico'),
        ('Fase 5: Validación (Sem 25-26)', 'Validación y documentación')
    ]
    
    tf.text = phases[0][0]
    tf.paragraphs[0].font.size = Pt(21)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    p = tf.add_paragraph()
    p.text = f'   → {phases[0][1]}'
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.space_after = Pt(10)
    
    for i in range(1, len(phases)):
        p = tf.add_paragraph()
        p.text = phases[i][0]
        p.font.size = Pt(21)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 51, 102)
        
        p = tf.add_paragraph()
        p.text = f'   → {phases[i][1]}'
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(10)
    
    # === SLIDE 11: CRONOGRAMA ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 248, 250)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = 'Cronograma (26 semanas)'
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(7), Inches(4))
    tf = content_box.text_frame
    tf.text = '📅 Distribución de Fases:'
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    
    phases = [
        '',
        '• Formulación: 8 semanas',
        '• Implementación: 8 semanas',
        '• Datos: 4 semanas',
        '• Experimentación: 4 semanas',
        '• Validación: 2 semanas'
    ]
    
    for phase in phases:
        p = tf.add_paragraph()
        p.text = phase
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(30, 30, 30)
        p.space_after = Pt(12)
    
    # === SLIDE 12: RESULTADOS ESPERADOS ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 248, 250)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = 'Resultados Esperados'
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5))
    tf = content_box.text_frame
    
    sections = [
        ('Productos Académicos', True),
        ('• Modelo matemático DLBP validado', False),
        ('• Artículo científico en revista indexada', False),
        ('• Dataset público para la comunidad', False),
        ('', False),
        ('Productos Tecnológicos', True),
        ('• Prototipo de software de planificación', False),
        ('• Código abierto en GitHub', False),
        ('', False),
        ('Impacto Esperado', True),
        ('• Reducción de costos (Ref: ~8.6%)', False),
        ('• Mejora en competitividad del sector', False)
    ]
    
    tf.text = sections[0][0]
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0, 100, 0)
    
    for text, is_header in sections[1:]:
        p = tf.add_paragraph()
        p.text = text
        if is_header:
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 100, 0)
        else:
            p.font.size = Pt(19)
            p.font.color.rgb = RGBColor(30, 30, 30)
        p.space_after = Pt(6)
    
    # === SLIDE 13: CONTRIBUCIONES ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 248, 250)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = 'Contribuciones Esperadas'
    title_frame.paragraphs[0].font.size = Pt(38)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5))
    tf = content_box.text_frame
    
    sections = [
        ('A la Ciencia', True),
        ('• Adaptación de modelos DLBP a coproductos', False),
        ('• Comparación rigurosa de metaheurísticas', False),
        ('', False),
        ('A la Industria', True),
        ('• Herramienta práctica de planificación', False),
        ('• Reducción cuantificable de costos', False),
        ('', False),
        ('A la Formación', True),
        ('• Competencias en optimización avanzada', False),
        ('• Vínculos academia-industria', False)
    ]
    
    tf.text = sections[0][0]
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(70, 130, 180)
    
    for text, is_header in sections[1:]:
        p = tf.add_paragraph()
        p.text = text
        if is_header:
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(70, 130, 180)
        else:
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(30, 30, 30)
        p.space_after = Pt(8)
    
    # === SLIDE 14: REFERENCIAS ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 248, 250)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = 'Referencias Clave'
    title_frame.paragraphs[0].font.size = Pt(38)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5))
    tf = content_box.text_frame
    
    refs = [
        'Becker, C., & Scholl, A. (1998). Survey on problems and methods in generalized assembly line balancing.',
        '',
        'Solano-Blanco, A. L., et al. (2022). Integrated planning decisions in the broiler chicken supply chain.',
        '',
        'Marel. (s.f.). El equilibrio de la carcasa: un reto clave para la industria avícola.',
        '',
        'Güngör, A., & Gupta, S. M. (2021). Disassembly scheduling: Literature review.',
        '',
        '[Bibliografía completa en protocolo de investigación]'
    ]
    
    tf.text = refs[0]
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = RGBColor(30, 30, 30)
    
    for i, ref in enumerate(refs[1:], 1):
        p = tf.add_paragraph()
        p.text = ref
        if ref.startswith('['):
            p.font.italic = True
            p.font.size = Pt(18)
        else:
            p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(30, 30, 30)
        p.space_after = Pt(4)
    
    # === SLIDE 15: CIERRE ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 248, 255)
    
    closing_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2.5))
    closing_frame = closing_box.text_frame
    closing_frame.text = '¡Gracias por su atención!'
    closing_frame.paragraphs[0].font.size = Pt(46)
    closing_frame.paragraphs[0].font.bold = True
    closing_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    closing_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    p = closing_frame.add_paragraph()
    p.text = ''
    p = closing_frame.add_paragraph()
    p.text = '¿Preguntas?'
    p.font.size = Pt(36)
    p.font.color.rgb = RGBColor(70, 130, 180)
    p.alignment = PP_ALIGN.CENTER
    
    contact_box = slide.shapes.add_textbox(Inches(1), Inches(5.8), Inches(8), Inches(1))
    contact_frame = contact_box.text_frame
    contact_frame.text = 'Daniel Castañeda'
    p = contact_frame.add_paragraph()
    p.text = 'Universidad Tecnológica de Pereira'
    for paragraph in contact_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.color.rgb = RGBColor(64, 64, 64)
    
    # Guardar presentación
    filename = 'presentacion_oral_protocolo.pptx'
    prs.save(filename)
    print(f'✅ Presentación creada exitosamente: {filename}')
    print(f'📊 Total de diapositivas: {len(prs.slides)}')
    return filename

if __name__ == '__main__':
    crear_presentacion()
